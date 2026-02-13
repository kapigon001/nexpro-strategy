"""共通ヘルパー関数 - McKinsey風スライドレイアウト"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# === カラーパレット ===
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
DARK_NAVY = RGBColor(0x12, 0x1D, 0x33)
GOLD = RGBColor(0xC8, 0xA9, 0x51)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x27, 0xAE, 0x60)
BLUE = RGBColor(0x34, 0x98, 0xDB)
LIGHT_GRAY = RGBColor(0xF2, 0xF3, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MED_GRAY = RGBColor(0x95, 0xA5, 0xA6)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)

FONT_NAME = "IPAGothic"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=12,
                bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                font_name=FONT_NAME, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.line_spacing = Pt(int(font_size * line_spacing))
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=11, color=DARK_TEXT, bold=False,
                          bullet=False, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        prefix = "• " if bullet else ""
        p.text = prefix + line
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.line_spacing = Pt(int(font_size * line_spacing))
    return txBox


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape


def set_shape_text(shape, text, font_size=11, color=WHITE, bold=False,
                   alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = alignment
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    return tf


def set_shape_multiline(shape, lines, font_size=10, color=WHITE, bold=False,
                        alignment=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.alignment = alignment


def add_footer(slide, page_num):
    # Page number
    add_textbox(slide, Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3),
                str(page_num), font_size=8, color=MED_GRAY,
                alignment=PP_ALIGN.RIGHT)
    # Confidential
    add_textbox(slide, Inches(0.3), Inches(7.1), Inches(1.5), Inches(0.3),
                "Confidential", font_size=8, color=MED_GRAY)


def add_header_bar(slide, title, message=None):
    """McKinseyスタイルのヘッダーバー"""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
                title, font_size=24, bold=True, color=WHITE)
    if message:
        add_rect(slide, Inches(0.3), Inches(1.2), Inches(12.7), Inches(0.5),
                 LIGHT_GRAY, border_color=GOLD)
        add_textbox(slide, Inches(0.5), Inches(1.22), Inches(12.3), Inches(0.45),
                    message, font_size=13, bold=True, color=NAVY,
                    alignment=PP_ALIGN.LEFT)


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_table(slide, rows, cols, data, left, top, width, height,
              header_color=NAVY, header_text_color=WHITE,
              highlight_rows=None, highlight_color=None):
    """テーブルを追加。data = [[row0col0, row0col1, ...], ...]"""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # 列幅を均等に
    col_w = int(width / cols)
    for i in range(cols):
        table.columns[i].width = col_w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if r < len(data) and c < len(data[r]) else ""

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.name = FONT_NAME

                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = header_text_color
                else:
                    paragraph.font.color.rgb = DARK_TEXT

            # ヘッダー行の背景
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif highlight_rows and r in highlight_rows and highlight_color:
                cell.fill.solid()
                cell.fill.fore_color.rgb = highlight_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table_shape


# === セクションディバイダー ===
def make_section_divider(prs, section_num, title, subtitle, page_num):
    slide = add_slide(prs)
    set_bg(slide, NAVY)
    add_textbox(slide, Inches(1), Inches(2.0), Inches(3), Inches(2),
                section_num, font_size=72, bold=True, color=GOLD)
    add_textbox(slide, Inches(1), Inches(3.8), Inches(10), Inches(1),
                title, font_size=36, bold=True, color=WHITE)
    add_textbox(slide, Inches(1), Inches(5.0), Inches(10), Inches(0.6),
                subtitle, font_size=16, color=MED_GRAY)
    add_footer(slide, page_num)
    return slide
