#!/usr/bin/env python3
# pip install python-pptx
"""
generate_presentation.py
========================
ネクプロ 全社戦略提案 2026 — 12枚 .pptx 自動生成スクリプト

Usage:
    python generate_presentation.py  →  presentation.pptx が生成される
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ─────────────────────────────────────────────────────────────
# ① 定数ブロック（微調整はここだけ）
# ─────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.333)
SLIDE_H  = Inches(7.5)

# カラー
C_ORANGE      = RGBColor(255, 107,   0)
C_WHITE       = RGBColor(255, 255, 255)
C_BLACK       = RGBColor(  0,   0,   0)
C_DARK_GRAY   = RGBColor( 31,  41,  55)
C_MID_GRAY    = RGBColor(107, 114, 128)
C_LIGHT_GRAY  = RGBColor(243, 244, 246)
C_ORANGE_BG   = RGBColor(255, 237, 213)
C_PHASE1      = RGBColor(254, 243, 199)   # 黄：Phase 1
C_PHASE2      = RGBColor(220, 252, 231)   # 緑：Phase 2
C_PHASE3      = RGBColor(219, 234, 254)   # 青：Phase 3

# フォント
FONT = 'Meiryo'

# 共通座標
MARGIN_L    = Inches(0.55)
ACCENT_W    = Inches(0.06)
LEAD_TOP    = Inches(0.28)
LEAD_H      = Inches(0.70)
CONTENT_TOP = Inches(1.12)
CONTENT_W   = Inches(12.23)
FOOTER_TOP  = Inches(7.15)

# ─────────────────────────────────────────────────────────────
# ② Helper 関数
# ─────────────────────────────────────────────────────────────

def _blank(prs):
    """完全空白スライドを追加して返す"""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, l, t, w, h, fill_rgb, border_rgb=None):
    """塗りつぶし矩形を追加して shape を返す"""
    shp = slide.shapes.add_shape(1, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    if border_rgb:
        shp.line.color.rgb = border_rgb
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    return shp


def _tb(slide, text, l, t, w, h,
        size=12, bold=False, color=C_BLACK,
        align=PP_ALIGN.LEFT, italic=False):
    """テキストボックスを追加して txBox を返す"""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = text
    run.font.name      = FONT
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    return txb


def _cell_style(cell, bg=None, fg=C_BLACK, bold=False, size=10,
                align=PP_ALIGN.LEFT, text=None):
    """テーブルセルのスタイルを設定"""
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg

    tf = cell.text_frame
    p  = tf.paragraphs[0]

    # 既存 run を削除
    for r_elem in p._p.findall(qn('a:r')):
        p._p.remove(r_elem)

    p.alignment = align
    if text is not None:
        run = p.add_run()
        run.text           = text
        run.font.name      = FONT
        run.font.size      = Pt(size)
        run.font.bold      = bold
        run.font.color.rgb = fg


def _table(slide, data, col_widths, left, top,
           row_h=Inches(0.40),
           hdr_bg=C_DARK_GRAY, hdr_fg=C_WHITE,
           highlight_col=None,
           phase_rows=None):
    """
    スタイル済みテーブルを追加
      data          : List[List[str]]  ※ data[0] がヘッダー行
      highlight_col : int              ※ 該当列を C_ORANGE_BG + C_ORANGE で強調
      phase_rows    : {row_idx: bg_rgb}
    """
    rows = len(data)
    cols = len(data[0])
    total_w = sum(col_widths)

    tshp = slide.shapes.add_table(rows, cols, left, top, total_w, row_h * rows)
    tbl  = tshp.table

    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = cw

    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            cell = tbl.cell(r, c)
            is_hdr = (r == 0)
            is_hl  = (highlight_col is not None and c == highlight_col and not is_hdr)
            is_alt = (r % 2 == 0 and r > 0)

            if is_hdr:
                bg, fg, bold = hdr_bg, hdr_fg, True
            elif phase_rows and r in phase_rows:
                bg, fg, bold = phase_rows[r], C_BLACK, False
            elif is_hl:
                bg, fg, bold = C_ORANGE_BG, C_ORANGE, True
            elif is_alt:
                bg, fg, bold = C_LIGHT_GRAY, C_BLACK, False
            else:
                bg, fg, bold = C_WHITE, C_BLACK, False

            _cell_style(cell, bg=bg, fg=fg, bold=bold, size=10, text=val)

    return tshp


def _accent_bar(slide):
    """左端オレンジ縦バー"""
    _rect(slide, Inches(0), Inches(0), ACCENT_W, SLIDE_H, C_ORANGE)


def _lead_bar(slide, message):
    """リードメッセージバー（薄グレー帯 + オレンジ下線）"""
    _rect(slide, MARGIN_L, LEAD_TOP, CONTENT_W, LEAD_H, C_LIGHT_GRAY)
    _rect(slide, MARGIN_L, LEAD_TOP + LEAD_H - Inches(0.04),
          CONTENT_W, Inches(0.04), C_ORANGE)
    _tb(slide, message,
        MARGIN_L + Inches(0.15), LEAD_TOP + Inches(0.10),
        CONTENT_W - Inches(0.3), LEAD_H - Inches(0.16),
        size=13, bold=True, color=C_DARK_GRAY)


def _frame_label(slide, label):
    """右上フレームラベル（CORE / WHY / WHAT / HOW）"""
    _tb(slide, label,
        Inches(11.5), Inches(0.06), Inches(1.7), Inches(0.28),
        size=11, bold=True, color=C_ORANGE, align=PP_ALIGN.RIGHT)


def _footer(slide, page_num):
    """フッター"""
    _tb(slide, f'Confidential  |  {page_num} / 12',
        Inches(9.0), FOOTER_TOP, Inches(4.2), Inches(0.30),
        size=9, color=C_MID_GRAY, align=PP_ALIGN.RIGHT)


def _base(slide, lead_msg, frame, page):
    """Pattern B 共通レイヤー（白背景スライド用）"""
    _accent_bar(slide)
    _lead_bar(slide, lead_msg)
    _frame_label(slide, frame)
    _footer(slide, page)


def _slide_title_tb(slide, text):
    """スライドタイトルテキストボックス（Pattern B 共通）"""
    _tb(slide, text,
        MARGIN_L, CONTENT_TOP, CONTENT_W, Inches(0.48),
        size=18, bold=True, color=C_DARK_GRAY)


# ─────────────────────────────────────────────────────────────
# ③ スライド生成関数（12本）
# ─────────────────────────────────────────────────────────────

def slide_01_title(prs):
    s = _blank(prs)
    _rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, C_DARK_GRAY)
    _accent_bar(s)

    _tb(s, 'ネクプロ 全社戦略提案 2026',
        Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.1),
        size=40, bold=True, color=C_WHITE)

    _tb(s, 'ウェビナーツールから、B2B成長インフラへ。',
        Inches(0.8), Inches(2.55), Inches(11.5), Inches(0.75),
        size=22, color=C_ORANGE)

    _rect(s, Inches(0.8), Inches(3.45), Inches(8.0), Inches(0.03), C_MID_GRAY)

    _tb(s, 'アジェンダ',
        Inches(0.8), Inches(3.65), Inches(11.0), Inches(0.45),
        size=13, bold=True, color=C_LIGHT_GRAY)

    for i, item in enumerate([
        '1.  なぜ今、転換が必要か（外部環境・競合・顧客変化）',
        '2.  どこを目指すか（ビジョン・戦略・ロードマップ）',
        '3.  今日決めること（資源配分・ポジショニング・組織）',
    ]):
        _tb(s, item,
            Inches(0.9), Inches(4.2) + i * Inches(0.55), Inches(11.0), Inches(0.5),
            size=14, color=C_WHITE)

    _tb(s, '2026年3月  |  Confidential',
        Inches(0.8), Inches(6.9), Inches(6.0), Inches(0.38),
        size=10, color=C_MID_GRAY)


def slide_02_exec_summary(prs):
    s = _blank(prs)
    _base(s, '配信ツールに留まれば死ぬ。データ基盤に転換すれば勝てる。', 'CORE', 2)
    _slide_title_tb(s, 'エグゼクティブサマリー')

    # 2×2 マトリクス
    cx = MARGIN_L
    cy = CONTENT_TOP + Inches(0.55)
    cw = Inches(2.75)
    ch = Inches(1.05)
    cells_info = [
        ('危険地帯\n（成長性：高 / 競争力：低）', C_LIGHT_GRAY, C_MID_GRAY, False),
        ('★ 目標\nエンゲージメント基盤',          C_ORANGE_BG,  C_ORANGE,   True),
        ('→ 現在地\n（配信ツール）',               C_LIGHT_GRAY, C_MID_GRAY, False),
        ('維持地帯',                               C_LIGHT_GRAY, C_MID_GRAY, False),
    ]
    positions = [(0, 0), (0, 1), (1, 0), (1, 1)]
    for (row, col), (txt, bg, fg, bold) in zip(positions, cells_info):
        x = cx + col * (cw + Inches(0.08))
        y = cy + row * (ch + Inches(0.06))
        _rect(s, x, y, cw, ch, bg, border_rgb=C_MID_GRAY)
        _tb(s, txt, x + Inches(0.1), y + Inches(0.12),
            cw - Inches(0.2), ch - Inches(0.24),
            size=10, bold=bold, color=fg)

    _tb(s, '成長性 ↑ / 競争力 →',
        cx, cy + ch * 2 + Inches(0.18), Inches(3.5), Inches(0.3),
        size=8, color=C_MID_GRAY)

    # サマリーテーブル
    data = [
        ['軸',       '現状（FY24）',       '目標（FY27）'],
        ['事業定義', 'ウェビナー配信ツール', 'B2Bエンゲージメント基盤'],
        ['ARR',      '¥512M',              '¥1,382M'],
        ['MRR解約率', '1.7%/月',           '1.0%/月'],
        ['ARPA',     '¥150K/月',           '¥204K/月'],
        ['新収益柱', '¥0',                 '¥378M'],
    ]
    _table(s, data, [Inches(1.8), Inches(2.3), Inches(2.3)],
           MARGIN_L + Inches(6.1), CONTENT_TOP + Inches(0.55))

    # 本日の意思決定
    dy = CONTENT_TOP + Inches(3.05)
    _tb(s, '本日の意思決定（3点）:',
        MARGIN_L, dy, CONTENT_W, Inches(0.38),
        size=12, bold=True, color=C_DARK_GRAY)
    for i, d in enumerate([
        '① 資源配分：プロダクト 45% / GTM 35% / 組織 20%',
        '② ポジショニング転換：「ウェビナーツール」→「エンゲージメント基盤」',
        '③ 組織再編：PMM / RevOps 兼務設置 + KPI オーナー制度',
    ]):
        _tb(s, d, MARGIN_L + Inches(0.1), dy + Inches(0.42) + i * Inches(0.42),
            CONTENT_W - Inches(0.2), Inches(0.4), size=11, color=C_DARK_GRAY)


def slide_03_saas_shift(prs):
    s = _blank(prs)
    _base(s, 'AIはUIを圧縮する。データを持つ者が勝つ。', 'WHY', 3)
    _slide_title_tb(s, '外部環境：SaaSの構造変化')

    layers = [
        ('AIエージェント（実行層）', C_ORANGE_BG,             C_ORANGE,             True,  '↑ 拡大'),
        ('SaaS UI（圧縮対象）',      C_LIGHT_GRAY,            C_MID_GRAY,           False, '↓ コモディティ化'),
        ('データ基盤（価値増大）',   RGBColor(219, 234, 254), RGBColor(30, 64, 175), True, '↑ 勝ち残り'),
    ]
    blk_l = MARGIN_L
    blk_t = CONTENT_TOP + Inches(0.55)
    blk_w = Inches(5.8)
    blk_h = Inches(1.0)
    for i, (lbl, bg, fg, bold, trend) in enumerate(layers):
        y = blk_t + i * (blk_h + Inches(0.05))
        _rect(s, blk_l, y, blk_w, blk_h, bg, border_rgb=C_MID_GRAY)
        _tb(s, lbl,   blk_l + Inches(0.2),  y + Inches(0.25), Inches(3.5), Inches(0.55),
            size=13, bold=bold, color=fg)
        _tb(s, trend, blk_l + Inches(4.0),  y + Inches(0.25), Inches(1.6), Inches(0.55),
            size=11, bold=bold, color=fg)

    data = [
        ['ファクト',                   '数値'],
        ['SaaSstock 市場消失',          '$2,850億（2026年2月）'],
        ['バーティカルSaaS 年初来下落', '−43%'],
        ['AIエージェントPJ中止予測',   '40%以上（2027年末 / Gartner）'],
        ['日本ウェビナーSaaS AI統合遅れ', '12〜18 ヶ月'],
    ]
    _table(s, data, [Inches(3.6), Inches(2.5)],
           MARGIN_L + Inches(6.2), CONTENT_TOP + Inches(0.55))

    nt = CONTENT_TOP + Inches(3.85)
    _tb(s, 'ネクプロへの示唆:',
        MARGIN_L, nt, CONTENT_W, Inches(0.38), size=12, bold=True, color=C_DARK_GRAY)
    for i, txt in enumerate([
        '• 配信機能 = AIが蚕食する戦場（battleground）',
        '• エンゲージメントデータ = AI活用で価値増大（gold mine）',
    ]):
        _tb(s, txt, MARGIN_L + Inches(0.1), nt + Inches(0.42) + i * Inches(0.40),
            CONTENT_W - Inches(0.2), Inches(0.38), size=11, color=C_DARK_GRAY)


def slide_04_jtbd(prs):
    s = _blank(prs)
    _base(s, '顧客が買うのは「配信機能」ではなく「商談と売上」。', 'WHY', 4)
    _slide_title_tb(s, '顧客の本当の課題（JTBD）')

    q_t = CONTENT_TOP + Inches(0.55)
    _rect(s, MARGIN_L, q_t, CONTENT_W, Inches(0.85), C_LIGHT_GRAY)
    _rect(s, MARGIN_L, q_t, Inches(0.05), Inches(0.85), C_ORANGE)
    _tb(s, 'ウェビナーを「やらされている」企業。\n専任者なし、予算薄、ノウハウなし。でも捨てられない施策。',
        MARGIN_L + Inches(0.2), q_t + Inches(0.10),
        CONTENT_W - Inches(0.3), Inches(0.70),
        size=12, italic=True, color=C_DARK_GRAY)

    half = CONTENT_W // 2
    mid  = MARGIN_L + half + Inches(0.08)
    pt   = CONTENT_TOP + Inches(1.55)
    _rect(s, MARGIN_L, pt, half - Inches(0.04), Inches(0.36), C_DARK_GRAY)
    _tb(s, 'Pain — 現状の痛み',
        MARGIN_L + Inches(0.1), pt + Inches(0.05),
        half - Inches(0.2), Inches(0.28), size=11, bold=True, color=C_WHITE)
    _rect(s, mid, pt, half - Inches(0.04), Inches(0.36), C_ORANGE)
    _tb(s, 'Gain — 目指す状態',
        mid + Inches(0.1), pt + Inches(0.05),
        half - Inches(0.2), Inches(0.28), size=11, bold=True, color=C_WHITE)

    pains = [
        '申込・視聴・商談のデータが繋がらない',
        'ROI説明できない → 予算取れない → 改善不可',
        '視聴52分（メール3秒の1,040倍）が未計測',
    ]
    gains = [
        '「なんとなく回る」→ 「成果の根拠がある」',
        '胸を張って予算を取りに行ける施策に転換',
        'ウェビナー = コストから「投資」へ変える',
    ]
    for i, (p_txt, g_txt) in enumerate(zip(pains, gains)):
        y = pt + Inches(0.40) + i * Inches(0.42)
        _tb(s, f'• {p_txt}', MARGIN_L + Inches(0.1), y,
            half - Inches(0.2), Inches(0.40), size=10, color=C_DARK_GRAY)
        _tb(s, f'• {g_txt}', mid + Inches(0.1), y,
            half - Inches(0.2), Inches(0.40), size=10, color=C_DARK_GRAY)

    data = [
        ['部門',   '本当に必要なもの'],
        ['マーケ', '良質リードの証明（MA連携・スコアリング）'],
        ['営業',   '温度感の見える化（視聴ログ → 商談優先順位）'],
        ['経営',   '施策ROIの可視化（投資対効果の数字）'],
    ]
    _table(s, data, [Inches(1.5), Inches(10.6)],
           MARGIN_L, CONTENT_TOP + Inches(3.3), row_h=Inches(0.38))


def slide_05_competitive(prs):
    s = _blank(prs)
    _base(s, '3つのモートで「配信戦争」から離脱する。', 'WHY', 5)
    _slide_title_tb(s, '競合マップと自社の立ち位置')

    ml = MARGIN_L
    mt = CONTENT_TOP + Inches(0.55)
    mw = Inches(5.8)
    mh = Inches(3.4)
    _rect(s, ml, mt, mw, mh, C_LIGHT_GRAY)

    _tb(s, '日本企業適合性 高 ↑', ml, mt - Inches(0.28), Inches(3.0), Inches(0.28),
        size=9, color=C_MID_GRAY)
    _tb(s, 'データ活用高度性 高 →', ml, mt + mh + Inches(0.04), Inches(4.0), Inches(0.25),
        size=9, color=C_MID_GRAY)

    plots = [
        ('EventHub / bizibl', 0.20, 0.20, C_MID_GRAY, False),
        ('Zoom / Teams',      0.12, 0.72, C_MID_GRAY, False),
        ('ON24',              0.60, 0.58, C_MID_GRAY, False),
        ('ネクプロ（現在）',  0.55, 0.28, C_ORANGE,   True),
        ('ネクプロ（目標）',  0.82, 0.08, C_ORANGE,   True),
    ]
    for lbl, rx, ry, color, bold in plots:
        dx = ml + int(mw * rx) - Inches(0.13)
        dy = mt + int(mh * ry) - Inches(0.13)
        _rect(s, dx, dy, Inches(0.26), Inches(0.26), color)
        _tb(s, lbl, dx - Inches(0.4), dy - Inches(0.30),
            Inches(2.0), Inches(0.28), size=9, bold=bold, color=color)

    data = [
        ['モート',               '内容'],
        ['Japan Enterprise Ops',  '日本特有の業務運用（稟議・インボイス・監査）への深度適合'],
        ['Salesforce Native ROI', 'エンゲージメント×売上データを1画面で可視化'],
        ['Dual-ID Architecture',  '顧客データ資産を保護しながら自社ID資産を積み上げる'],
    ]
    _table(s, data, [Inches(2.5), Inches(5.9)],
           MARGIN_L + Inches(6.1), mt)

    nt = mt + mh + Inches(0.2)
    _rect(s, MARGIN_L, nt, CONTENT_W, Inches(0.55), C_ORANGE_BG, border_rgb=C_ORANGE)
    _tb(s, '⚡  ON24のCvent買収（2025年12月）で日本市場の優先度が低下。今が先行する最大の時間的窓。',
        MARGIN_L + Inches(0.15), nt + Inches(0.10),
        CONTENT_W - Inches(0.3), Inches(0.40),
        size=11, bold=True, color=C_ORANGE)


def slide_06_strategy_3pillar(prs):
    s = _blank(prs)
    _base(s, '「作る・集める・売る」を統合するプラットフォームへ。', 'CORE', 6)
    _slide_title_tb(s, '事業戦略：Global B2B Growth Platform')

    tiles = [
        ('AI Production\n（作る）',
         'AI字幕・要約・翻訳\nコンテンツ再資産化',
         '制作コスト・リードタイム削減'),
        ('Open Media\n（集める）',
         '公開メディア基盤\nSEO・レコメンドエンジン',
         '新規流入の継続的獲得'),
        ('Revenue Intelligence\n（売る）',
         'エンゲージメントスコア\nROIダッシュボード',
         '施策 → 売上の可視化'),
    ]
    tw = Inches(3.9)
    th = Inches(2.9)
    tt = CONTENT_TOP + Inches(0.55)
    gap = Inches(0.13)

    for i, (title, func, val) in enumerate(tiles):
        tl = MARGIN_L + i * (tw + gap)
        _rect(s, tl, tt, tw, th, C_WHITE, border_rgb=C_MID_GRAY)
        _rect(s, tl, tt, tw, Inches(0.65), C_ORANGE)
        _tb(s, title, tl + Inches(0.1), tt + Inches(0.06),
            tw - Inches(0.2), Inches(0.58), size=13, bold=True, color=C_WHITE)
        _tb(s, func,  tl + Inches(0.15), tt + Inches(0.78),
            tw - Inches(0.3), Inches(0.95), size=11, color=C_DARK_GRAY)
        _tb(s, val,   tl + Inches(0.15), tt + Inches(1.85),
            tw - Inches(0.3), Inches(0.8),  size=11, bold=True, color=C_ORANGE)

    _tb(s, 'SaaS × BPaaS × コミュニティ:  プロダクトだけでは届かない「ノウハウと人材の不足」をBPaaSで補完。コミュニティが実践知の循環を生む。',
        MARGIN_L, tt + th + Inches(0.22), CONTENT_W, Inches(0.6),
        size=11, color=C_DARK_GRAY)


def slide_07_flywheel(prs):
    s = _blank(prs)
    _base(s, '利用がそのまま集客とデータ蓄積になる自己強化ループ。', 'WHAT', 7)
    _slide_title_tb(s, 'フライホイール設計')

    fw_steps = [
        ('AI制作（Input）',              C_ORANGE,    C_WHITE,    True),
        ('コンテンツ公開',               C_LIGHT_GRAY, C_DARK_GRAY, False),
        ('検索・自然流入（Traffic）',    C_LIGHT_GRAY, C_DARK_GRAY, False),
        ('エンゲージメント（Insight）',  C_LIGHT_GRAY, C_DARK_GRAY, False),
        ('ROI可視化・商談創出',          C_ORANGE_BG, C_ORANGE,   True),
        ('顧客成功・事例蓄積',           C_LIGHT_GRAY, C_DARK_GRAY, False),
    ]
    ew = Inches(2.6)
    eh = Inches(0.68)
    ex = [MARGIN_L + Inches(0.1), MARGIN_L + Inches(0.1) + ew + Inches(0.55)]
    ey = [CONTENT_TOP + Inches(0.6), CONTENT_TOP + Inches(1.45), CONTENT_TOP + Inches(2.3)]

    coords = [(0,0),(1,0),(1,1),(1,2),(0,2),(0,1)]
    for step_i, (col, row) in enumerate(coords):
        x = ex[col]; y = ey[row]
        lbl, bg, fg, bold = fw_steps[step_i]
        _rect(s, x, y, ew, eh, bg, border_rgb=C_MID_GRAY)
        _tb(s, lbl, x + Inches(0.1), y + Inches(0.12),
            ew - Inches(0.2), eh - Inches(0.24), size=10, bold=bold, color=fg)

    _tb(s, '→ → ↓ ← ← ↑  （自己強化ループ）',
        MARGIN_L, CONTENT_TOP + Inches(3.15), Inches(6.5), Inches(0.35),
        size=10, color=C_MID_GRAY, align=PP_ALIGN.CENTER)

    data = [
        ['先行KPI',                       '遅行KPI'],
        ['コンテンツ公開本数 / 自然流入率', 'MQL化率 / 受注寄与率'],
        ['エンゲージメントスコア平均',      'ARPA / NRR'],
        ['オンボーディング完了率',          '月次解約率'],
    ]
    _table(s, data, [Inches(3.1), Inches(3.1)],
           MARGIN_L + Inches(6.5), CONTENT_TOP + Inches(0.6))


def slide_08_moat_japan(prs):
    s = _blank(prs)
    _base(s, '日本企業運用の複雑性が最大の参入障壁になる。', 'WHY', 8)
    _slide_title_tb(s, 'モート①：Japan Enterprise Ops')

    data = [
        ['要件',               'ネクプロ', 'Zoom', 'ON24', 'Teams'],
        ['日本語運用ネイティブ', '◎',      '△',   '△',   '○'],
        ['稟議・インボイス対応', '◎',      '✗',   '✗',   '△'],
        ['ISMS認証',            '◎',      '○',   '○',   '◎'],
        ['伴走型CS（日本語）',  '◎',      '✗',   '✗',   '△'],
        ['官公庁・医療対応',    '◎',      '△',   '✗',   '○'],
    ]
    _table(s, data, [Inches(3.5), Inches(1.8), Inches(1.5), Inches(1.5), Inches(1.5)],
           MARGIN_L, CONTENT_TOP + Inches(0.55), highlight_col=1)

    nt = CONTENT_TOP + Inches(3.35)
    _tb(s, '日本市場固有の参入コスト（グローバル勢の再現困難な理由）:',
        MARGIN_L, nt, CONTENT_W, Inches(0.38), size=12, bold=True, color=C_DARK_GRAY)
    for i, f in enumerate([
        'SaaS浸透率：日本 4%（米国 15〜18%）→ 導入支援・伴走CSが必須',
        'デジタル人材不足：2026年までに 230万人不足 → 運用代行ニーズ急増',
        '言語の壁：日本語バイヤーの 72% が母国語コンテンツ・サポートを選好',
    ]):
        _tb(s, f'• {f}',
            MARGIN_L + Inches(0.1), nt + Inches(0.42) + i * Inches(0.42),
            CONTENT_W - Inches(0.2), Inches(0.40), size=11, color=C_DARK_GRAY)


def slide_09_moat_salesforce(prs):
    s = _blank(prs)
    _base(s, 'エンゲージメント×売上データを1画面で見せるのはネクプロだけ。', 'WHY', 9)
    _slide_title_tb(s, 'モート②：Salesforce Native ROI')

    _tb(s, 'Salesforce 内で完結する分析:',
        MARGIN_L, CONTENT_TOP + Inches(0.55), CONTENT_W, Inches(0.38),
        size=12, bold=True, color=C_DARK_GRAY)
    for i, a in enumerate([
        'コンテンツ接触履歴（ウェビナー参加・視聴ログ）',
        '商談化・受注との相関（エンゲージメントスコア × Salesforce 商談）',
        '次アクション提案（Agentforce との連携）',
    ]):
        _tb(s, f'• {a}',
            MARGIN_L + Inches(0.1), CONTENT_TOP + Inches(1.02) + i * Inches(0.40),
            CONTENT_W - Inches(0.2), Inches(0.38), size=11, color=C_DARK_GRAY)

    data = [
        ['ツール',   '配信データ', 'CRM連携',                'ROI可視化'],
        ['ネクプロ', '◎',          '◎ Salesforceネイティブ', '◎（ロードマップ）'],
        ['Zoom',     '◎',          '○ API連携',              '✗'],
        ['MA単体',   '✗',          '◎',                      '△'],
        ['ON24',     '○',          '◎ Eloqua/Marketo',       '◎ ACE AI'],
    ]
    _table(s, data, [Inches(1.8), Inches(1.8), Inches(4.2), Inches(4.2)],
           MARGIN_L, CONTENT_TOP + Inches(2.2))

    nt = CONTENT_TOP + Inches(4.45)
    _rect(s, MARGIN_L, nt, CONTENT_W, Inches(0.52), C_ORANGE_BG, border_rgb=C_ORANGE)
    _tb(s, '⚡  Agentforce統合：国内最先行。ただし「序章」にすぎない — 本番はROIダッシュボード完成後。',
        MARGIN_L + Inches(0.15), nt + Inches(0.10),
        CONTENT_W - Inches(0.3), Inches(0.40), size=11, bold=True, color=C_ORANGE)


def slide_10_roadmap(prs):
    s = _blank(prs)
    _base(s, 'AI×データ×メディアの3フェーズで基盤を構築する。', 'WHAT', 10)
    _slide_title_tb(s, '製品ロードマップ（0〜18ヶ月）')

    data = [
        ['フェーズ', '期間',       'テーマ',       '主要リリース'],
        ['Phase 1',  '0〜6ヶ月',   '止血・種まき',
         'AI字幕/要約/翻訳、Salesforce埋め込みβ、エンゲージメントスコアMVP'],
        ['Phase 2',  '6〜12ヶ月',  '転換・仕込み',
         '公開メディア基盤、AIコンテンツ再資産化、レコメンドα'],
        ['Phase 3',  '12〜18ヶ月', '成長・回収',
         'ROIダッシュボード正式版、業種別BPaaSパッケージ'],
    ]
    _table(s, data, [Inches(1.2), Inches(1.5), Inches(2.0), Inches(7.4)],
           MARGIN_L, CONTENT_TOP + Inches(0.55),
           row_h=Inches(0.88), phase_rows={1: C_PHASE1, 2: C_PHASE2, 3: C_PHASE3})

    gt = CONTENT_TOP + Inches(4.2)
    _rect(s, MARGIN_L, gt, CONTENT_W, Inches(0.58), C_LIGHT_GRAY)
    _rect(s, MARGIN_L, gt, Inches(0.05), Inches(0.58), C_ORANGE)
    _tb(s, 'Gate Review（M6）: KPI達成（解約率1.3%以下 / エンゲージメントスコアMVP 5社導入）を条件に攻勢型（松案）へ移行判断。',
        MARGIN_L + Inches(0.2), gt + Inches(0.10),
        CONTENT_W - Inches(0.3), Inches(0.42), size=11, color=C_DARK_GRAY)


def slide_11_gtm_finance(prs):
    s = _blank(prs)
    _base(s, 'ARR ¥512M → ¥1,382M の道筋と3つの収益柱。', 'HOW', 11)
    _slide_title_tb(s, 'GTMと財務計画')

    data = [
        ['収益源',                          'FY24',     'FY25計画',  'FY26計画',   'FY27計画'],
        ['MRR',                             '¥287.8M',  '¥330M',    '¥417.4M',   '¥526.8M'],
        ['オプションサービス',              '¥227.1M',  '¥319.8M',  '¥494.5M',   '¥856M'],
        ['新収益柱（コンパウンド+営業DX）', '¥0',       '¥36.7M',   '¥138.8M',   '¥378M'],
        ['合計',                            '¥512.8M',  '¥649.8M',  '¥911.9M',   '¥1,382.8M'],
    ]
    _table(s, data, [Inches(3.7), Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.7)],
           MARGIN_L, CONTENT_TOP + Inches(0.55))

    data2 = [
        ['軌道',         '内容',                                    '対象'],
        ['スイッチ軌道', 'Zoomでは足りないと感じた企業へのリポジション提案', '既存市場の深耕'],
        ['BPaaS軌道',   '運用代行から始まりSaaSへ育てるハイタッチ型',      '新規市場の開拓'],
    ]
    _table(s, data2, [Inches(2.2), Inches(6.7), Inches(3.2)],
           MARGIN_L, CONTENT_TOP + Inches(3.22))

    _tb(s, '資源配分:  プロダクト 45%  /  GTM 35%  /  組織 20%',
        MARGIN_L, CONTENT_TOP + Inches(4.5), CONTENT_W, Inches(0.4),
        size=12, bold=True, color=C_DARK_GRAY)


def slide_12_decision(prs):
    s = _blank(prs)
    _base(s, '今日この場で決めなければ、「何もしない」という選択を自動的にしたことになる。',
          'HOW', 12)
    _slide_title_tb(s, '本日のお願い（3つの意思決定）')

    data = [
        ['#',    '決議内容',                                                                '推奨'],
        ['決議1', '資源配分の承認（プロダクト 45% / GTM 35% / 組織 20%）',                 '✅ 承認'],
        ['決議2', 'ポジショニング転換の承認（「ウェビナーツール」→「B2Bエンゲージメント基盤」）', '✅ 承認'],
        ['決議3', '組織再編の承認（PMM兼務・RevOps兼務・KPIオーナー制度・90日実行計画）',    '✅ 承認'],
    ]
    _table(s, data, [Inches(1.0), Inches(9.6), Inches(1.5)],
           MARGIN_L, CONTENT_TOP + Inches(0.55), row_h=Inches(0.62))

    _tb(s, '承認後の即時アクション（Next Steps）:',
        MARGIN_L, CONTENT_TOP + Inches(2.88), CONTENT_W, Inches(0.38),
        size=12, bold=True, color=C_DARK_GRAY)

    data2 = [
        ['アクション',                              '担当',          '期限'],
        ['KPIオーナー全員アサイン通知',             'CEO',           '1週間以内'],
        ['PMM兼務者の選定・辞令',                   'CEO',           '2週間以内'],
        ['エンゲージメントスコアMVP要件定義キックオフ', 'プロダクト責任者', '2週間以内'],
        ['業種別パッケージ作成チーム組成',           '営業責任者',     '2週間以内'],
        ['月次KPIレビュー初回日程設定',             'COO相当',        '1週間以内'],
    ]
    _table(s, data2, [Inches(5.5), Inches(2.8), Inches(3.8)],
           MARGIN_L, CONTENT_TOP + Inches(3.3), row_h=Inches(0.38))


# ─────────────────────────────────────────────────────────────
# ④ main
# ─────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_exec_summary(prs)
    slide_03_saas_shift(prs)
    slide_04_jtbd(prs)
    slide_05_competitive(prs)
    slide_06_strategy_3pillar(prs)
    slide_07_flywheel(prs)
    slide_08_moat_japan(prs)
    slide_09_moat_salesforce(prs)
    slide_10_roadmap(prs)
    slide_11_gtm_finance(prs)
    slide_12_decision(prs)

    out = 'presentation.pptx'
    prs.save(out)
    print(f'✅  Saved: {out}  ({len(prs.slides)} slides)')


if __name__ == '__main__':
    main()
